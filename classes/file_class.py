
from __future__ import annotations

import os
import re
import hashlib
import mimetypes
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Protocol

# ---------- Data models ----------

@dataclass(frozen=True)
class FileInfo:
    path: Path
    size: int
    mtime: float
    kind: str  # 'pdf' | 'pptx' | 'docx' | 'txt' | 'md' | 'unknown'

@dataclass
class ExtractResult:
    text: str
    metadata: Dict[str, str]
    chunks: Optional[List[str]] = None

# ---------- Extractor SPI ----------

class BaseExtractor(Protocol):
    kind: str
    def can_handle(self, info: FileInfo) -> bool: ...
    def extract_text(self, info: FileInfo) -> ExtractResult: ...

# ---------- Concrete extractors (optional deps; degrade gracefully) ----------

class PDFExtractor:
    kind = "pdf"
    def can_handle(self, info: FileInfo) -> bool:
        return info.kind == self.kind
    def extract_text(self, info: FileInfo) -> ExtractResult:
        text = ""
        meta = {"source": str(info.path), "type": "pdf"}
        try:
            from pypdf import PdfReader  # type: ignore
            reader = PdfReader(str(info.path))
            for page in reader.pages:
                t = page.extract_text() or ""
                text += t + "\n"
        except Exception:
            try:
                from pdfminer.high_level import extract_text as pm_extract  # type: ignore
                text = pm_extract(str(info.path))
            except Exception as e:
                raise RuntimeError(f"PDF extraction failed: {e}")
        return ExtractResult(text=text.strip(), metadata=meta)

class PPTXExtractor:
    kind = "pptx"
    def can_handle(self, info: FileInfo) -> bool:
        return info.kind == self.kind
    def extract_text(self, info: FileInfo) -> ExtractResult:
        try:
            from pptx import Presentation  # type: ignore
        except Exception as e:
            raise RuntimeError("python-pptx not installed") from e
        prs = Presentation(str(info.path))
        texts: List[str] = []
        for i, slide in enumerate(prs.slides):
            bits: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = (shape.text or "").strip()
                    if t:
                        bits.append(t)
            # notes (best-effort)
            try:
                if prs.has_notes_slide and getattr(slide, "notes_slide", None):
                    ntf = getattr(slide.notes_slide, "notes_text_frame", None)
                    if ntf and ntf.text:
                        bits.append(f"[Notes] {ntf.text.strip()}")
            except Exception:
                pass
            if bits:
                texts.append(f"[Slide {i+1}]\n" + "\n".join(bits))
        joined = "\n\n".join(texts)
        return ExtractResult(text=joined, metadata={"source": str(info.path), "type": "pptx", "slides": str(len(prs.slides))})

class DOCXExtractor:
    kind = "docx"
    def can_handle(self, info: FileInfo) -> bool:
        return info.kind == self.kind
    def extract_text(self, info: FileInfo) -> ExtractResult:
        try:
            import docx  # type: ignore
        except Exception as e:
            raise RuntimeError("python-docx not installed") from e
        d = docx.Document(str(info.path))
        blocks: List[str] = []
        for p in d.paragraphs:
            t = p.text.strip()
            if t:
                blocks.append(t)
        for table in d.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    blocks.append(" | ".join(cells))
        return ExtractResult(text="\n".join(blocks), metadata={"source": str(info.path), "type": "docx"})

class TextExtractor:
    kind = "txt"
    def can_handle(self, info: FileInfo) -> bool:
        return info.kind in {"txt", "md"}
    def extract_text(self, info: FileInfo) -> ExtractResult:
        encs = ["utf-8", "utf-8-sig", "cp949", "euc-kr", "latin1"]
        data = None
        last_err: Optional[Exception] = None
        for enc in encs:
            try:
                data = info.path.read_text(encoding=enc)
                break
            except Exception as e:
                last_err = e
        if data is None:
            raise RuntimeError(f"Text decode failed: {last_err}")
        return ExtractResult(text=data, metadata={"source": str(info.path), "type": info.kind})

DEFAULT_EXTRACTORS: List[BaseExtractor] = [
    PDFExtractor(), PPTXExtractor(), DOCXExtractor(), TextExtractor()
]

# ---------- Helpers ----------

_EXT_MAP = {".pdf":"pdf",".pptx":"pptx",".docx":"docx",".txt":"txt",".md":"md"}
_MIME_MAP = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "text/plain": "txt",
    "text/markdown": "md",
}

def _hash_file(path: Path, chunk: int = 1 << 20) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        while True:
            b = f.read(chunk)
            if not b:
                break
            h.update(b)
    return h.hexdigest()

# ---------- Multi-base File Manager ----------

class MultiBaseFileManager:
    """
    - allowed_bases: 여러 개의 허용 루트(샌드박스) 경로
    - 안전 경로 검사: 해석된 절대경로가 allowed_bases 중 하나의 하위여야 함
    - listdir/resolve_name/find_folder: LLM 보조용 탐색 유틸
    - extract_*: 파일 추출 (PDF/PPTX/DOCX/TXT)
    """
    def __init__(self, allowed_bases: List[Path], extractors: Optional[List[BaseExtractor]] = None):
        bases = [b.resolve() for b in (allowed_bases or []) if b]
        if not bases:
            raise ValueError("allowed_bases is empty")
        # remove duplicates
        seen = set()
        uniq: List[Path] = []
        for b in bases:
            s = str(b)
            if s not in seen:
                seen.add(s)
                uniq.append(b)
        self.allowed = uniq
        self.extractors = extractors or list(DEFAULT_EXTRACTORS)
        self._cache: Dict[str, ExtractResult] = {}

    # ---- safe resolve ----
    def _is_under_allowed(self, q: Path) -> bool:
        sq = str(q)
        for base in self.allowed:
            if sq.startswith(str(base)):
                return True
        return False

    def _safe_resolve(self, p: Path, *, default_base_index: int = 0) -> Path:
        q = p if p.is_absolute() else (self.allowed[default_base_index] / p)
        q = q.resolve()
        if not self._is_under_allowed(q):
            raise PermissionError(f"Path escapes allowed roots: {q}")
        return q

    # ---- kind & stat ----
    def detect_kind(self, path: Path) -> str:
        p = self._safe_resolve(path)
        ext_kind = _EXT_MAP.get(p.suffix.lower())
        if ext_kind:
            return ext_kind
        mime, _ = mimetypes.guess_type(str(p))
        return _MIME_MAP.get(mime, "unknown")

    def stat(self, path: Path) -> FileInfo:
        p = self._safe_resolve(path)
        if not p.exists() or not p.is_file():
            raise FileNotFoundError(str(p))
        return FileInfo(path=p, size=p.stat().st_size, mtime=p.stat().st_mtime, kind=self.detect_kind(p))

    # ---- list & fuzzy resolve ----
    def listdir(self, base_index: int = 0, rel: Path = Path("../../../Users/김호태/Downloads"), *, files_only: bool = False, max_items: int = 500) -> List[Dict[str, str]]:
        root = self.allowed[base_index]
        target = self._safe_resolve(root / rel if rel else root, default_base_index=base_index)
        if not target.exists() or not target.is_dir():
            raise FileNotFoundError(str(target))
        items: List[Dict[str, str]] = []
        for i, child in enumerate(sorted(target.iterdir(), key=lambda x: x.name.lower())):
            if i >= max_items:
                break
            try:
                is_file = child.is_file()
                if files_only and not is_file:
                    continue
                st = child.stat()
                items.append({
                    "name": child.name,
                    "type": "file" if is_file else "dir",
                    "size": str(st.st_size),
                    "mtime": str(st.st_mtime),
                    "ext": child.suffix.lower(),
                    "base": str(root),
                    "relpath": str(child.relative_to(root)),
                })
            except Exception:
                continue
        return items

    def resolve_name(self, user_name: str, base_index: int = 0, rel: Path = Path("../../../Users/김호태/Downloads"), *, cutoff: float = 0.6, topk: int = 5) -> Dict[str, object]:
        from difflib import SequenceMatcher
        listing = self.listdir(base_index=base_index, rel=rel)
        def score(a: str, b: str) -> float:
            return SequenceMatcher(None, a.lower(), b.lower()).ratio()
        scored = []
        for it in listing:
            s = score(user_name, it["name"]) if user_name else 0.0
            if s >= cutoff:
                scored.append({**it, "score": round(s, 3)})
        scored.sort(key=lambda x: x["score"], reverse=True)
        return {"input": user_name, "candidates": scored[:topk], "directory": str((self.allowed[base_index] / rel).resolve())}

    def find_folder(self, hint: str, *, max_depth: int = 2, topk: int = 5) -> Dict[str, object]:
        """Search all allowed bases for a folder whose name matches `hint` with fuzzy logic."""
        from difflib import SequenceMatcher

        def norm(s: str) -> str:
            return s.strip().lower()

        q = norm(hint)
        if not q:
            return {"found": False, "candidates": []}

        cands: List[Dict[str, object]] = []
        for base in self.allowed:
            for path, dirs, _files in os.walk(base):
                depth = len(Path(path).relative_to(base).parts)
                if depth > max_depth:
                    dirs[:] = []
                    continue
                for d in dirs:
                    name = d
                    p = Path(path) / d
                    s = SequenceMatcher(None, q, norm(name)).ratio()
                    if norm(name) == q:
                        s += 0.30
                    elif norm(name).startswith(q):
                        s += 0.15
                    elif q in norm(name):
                        s += 0.10
                    st = p.stat()
                    cands.append({
                        "base": str(base),
                        "relpath": str(p.relative_to(base)),
                        "name": name,
                        "mtime": st.st_mtime,
                        "score": round(min(1.0, s), 3),
                        "depth": depth,
                    })
        cands.sort(key=lambda x: (x["score"], -x["mtime"], -x["depth"]), reverse=True)
        if not cands:
            return {"found": False, "candidates": []}
        top = cands[:topk]
        if top and (top[0]["score"] >= 0.75) and (len(top) == 1 or top[0]["score"] - top[1]["score"] >= 0.1):
            return {"found": True, "best": top[0], "candidates": top}
        return {"found": "ambiguous", "candidates": top}

    # ---- extraction ----
    def _chunk_text(self, text: str, *, chunk_size: int = 1200, overlap: int = 150) -> List[str]:
        if not text:
            return []
        paras = re.split(r"\n{2,}", text)
        chunks: List[str] = []
        buf = ""
        for para in paras:
            if len(buf) + len(para) + 2 <= chunk_size:
                buf = (buf + "\n\n" + para).strip()
            else:
                if buf:
                    chunks.append(buf)
                if len(para) > chunk_size:
                    start = 0
                    while start < len(para):
                        end = start + chunk_size
                        chunks.append(para[start:end])
                        start = max(end - overlap, start + 1)
                    buf = ""
                else:
                    buf = para
        if buf:
            chunks.append(buf)
        return chunks

    def _extract_with(self, info: FileInfo) -> ExtractResult:
        cache_key = f"{_hash_file(info.path)}::{info.mtime}::{info.kind}"
        if cache_key in self._cache:
            res = self._cache[cache_key]
            return ExtractResult(text=res.text, metadata=dict(res.metadata), chunks=list(res.chunks) if res.chunks else None)
        extractor = next((ex for ex in self.extractors if ex.can_handle(info)), None)
        if extractor is None:
            raise ValueError(f"No extractor registered for kind={info.kind}")
        result = extractor.extract_text(info)
        result.metadata.setdefault("size", str(info.size))
        result.metadata.setdefault("mtime", str(info.mtime))
        result.metadata.setdefault("kind", info.kind)
        self._cache[cache_key] = result
        return result

    def extract_absolute(self, absolute_path: Path, *, do_chunk: bool = True, chunk_size: int = 1200, overlap: int = 150) -> ExtractResult:
        """Extract when you already have an absolute path under allowed roots."""
        p = absolute_path.resolve()
        if not self._is_under_allowed(p):
            raise PermissionError(f"Path escapes allowed roots: {p}")
        info = self.stat(p)
        res = self._extract_with(info)
        if do_chunk:
            res.chunks = self._chunk_text(res.text, chunk_size=chunk_size, overlap=overlap)
        return res

    def extract_from(self, base_index: int, rel: Path, *, do_chunk: bool = True, chunk_size: int = 1200, overlap: int = 150) -> ExtractResult:
        """Extract by (base_index, relative path)."""
        p = self._safe_resolve(self.allowed[base_index] / rel, default_base_index=base_index)
        info = self.stat(p)
        res = self._extract_with(info)
        if do_chunk:
            res.chunks = self._chunk_text(res.text, chunk_size=chunk_size, overlap=overlap)
        return res

    # ---- folder ops ----
    def make_dir(self, base_index: int, rel: Path, *, parents: bool = True, exist_ok: bool = True) -> Path:
        """
        지정한 베이스/상대경로에 폴더 생성.
        반환값: 생성(또는 존재)하는 디렉터리의 절대 경로
        """
        target = self._safe_resolve(self.allowed[base_index] / rel, default_base_index=base_index)
        target.mkdir(parents=parents, exist_ok=exist_ok)
        if not target.exists() or not target.is_dir():
            raise RuntimeError(f"Failed to create directory: {target}")
        return target

    def ensure_dir(self, base_index: int, rel: Path) -> Path:
        """폴더가 없으면 만들고, 있으면 그대로 반환."""
        return self.make_dir(base_index, rel, parents=True, exist_ok=True)

    # ---- move (cut & paste) ops ----
    def _unique_name(self, folder: Path, name: str) -> str:
        """
        folder 안에서 name이 겹치면 'name (1).ext', 'name (2).ext'... 형태로 비충돌 이름 생성
        """
        stem = Path(name).stem
        suffix = Path(name).suffix
        cand = name
        i = 1
        while (folder / cand).exists():
            cand = f"{stem} ({i}){suffix}"
            i += 1
        return cand

    def _move(
            self,
            src_base_index: int,
            src_rel: Path,
            dst_base_index: int,
            dst_folder_rel: Path,
            *,
            new_name: Optional[str] = None,
            overwrite: bool = False,
            dry_run: bool = False,
            _prepped_dst: Optional[Path] = None,  # move_many에서 미리 만든 폴더 재사용
    ) -> Dict[str, str]:
        src_abs = self._safe_resolve(self.allowed[src_base_index] / src_rel, default_base_index=src_base_index)
        if not src_abs.exists() or not src_abs.is_file():
            raise FileNotFoundError(f"Source file not found: {src_abs}")

        dst_folder_abs = _prepped_dst or self._safe_resolve(
            self.allowed[dst_base_index] / dst_folder_rel, default_base_index=dst_base_index
        )
        if not dst_folder_abs.exists() or not dst_folder_abs.is_dir():
            raise FileNotFoundError(f"Destination folder not found: {dst_folder_abs}")

        target_name = new_name or src_abs.name
        dst_abs = dst_folder_abs / target_name

        status = "moved"
        if dst_abs.exists():
            if overwrite:
                status = "overwritten"
            else:
                target_name = self._unique_name(dst_folder_abs, target_name)
                dst_abs = dst_folder_abs / target_name
                status = "renamed"

        if dry_run:
            return {"src": str(src_abs), "dst": str(dst_abs), "status": f"would_{status}"}

        if status == "overwritten" and dst_abs.exists():
            dst_abs.unlink()

        src_abs.replace(dst_abs)
        return {"src": str(src_abs), "dst": str(dst_abs), "status": status}

    def move_many(
            self,
            items: Union[Dict[str, object], List[Dict[str, object]]],
            dst_base_index: int,
            dst_folder_rel: Path,
            *,
            overwrite: bool = False,
            create_dst: bool = True,
            dry_run: bool = False,
            on_error: str = "continue",  # "continue" | "stop"
    ) -> Dict[str, object]:
        # 단일 dict도 허용 → 리스트로 정규화
        if isinstance(items, dict):
            items = [items]

        # 목적지 폴더 1회 준비
        dst_folder_abs = self._safe_resolve(self.allowed[dst_base_index] / dst_folder_rel,
                                            default_base_index=dst_base_index)
        if create_dst:
            dst_folder_abs.mkdir(parents=True, exist_ok=True)
        if not dst_folder_abs.exists() or not dst_folder_abs.is_dir():
            raise FileNotFoundError(f"Destination folder not found: {dst_folder_abs}")

        moved, errors = [], []

        for it in items:
            try:
                r = self._move(
                    src_base_index=int(it["src_base_index"]),
                    src_rel=Path(str(it["src_rel"])),
                    dst_base_index=dst_base_index,
                    dst_folder_rel=dst_folder_rel,
                    new_name=str(it.get("new_name")) if it.get("new_name") else None,
                    overwrite=overwrite,
                    dry_run=dry_run,
                    _prepped_dst=dst_folder_abs,  # 재사용
                )
                moved.append(r)
            except Exception as e:
                err = {"item": f'{it.get("src_base_index")}:{it.get("src_rel")}', "error": str(e)}
                errors.append(err)
                if on_error == "stop":
                    break

        return {"moved": moved, "errors": errors, "dst": str(dst_folder_abs)}